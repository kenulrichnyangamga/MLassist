from qdrant_client import models
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
from app.services.qdrant_service import get_qdrant_client
import uuid
import os

# Modell einmal laden, nicht bei jedem Aufruf
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

splitter = RecursiveCharacterTextSplitter(chunk_size=512, chunk_overlap=64)


def _split_seiten(seiten):
    """
    Nimmt eine Liste von (seitennummer, text) und gibt eine Liste von
    (chunk_text, seitennummer) zurück. So bleibt jede Seite ihrem Chunk zugeordnet.
    """
    chunk_liste = []
    for seitennummer, text in seiten:
        if not text or not text.strip():
            continue
        for chunk in splitter.split_text(text):
            chunk_liste.append((chunk, seitennummer))
    return chunk_liste


def extract_pdf(pdf_path):
    """PDF: Text pro Seite extrahieren, Seitennummer beibehalten."""
    doc = fitz.open(pdf_path)
    seiten = []
    for i, page in enumerate(doc):
        seiten.append((i + 1, page.get_text()))  # Seiten ab 1 zählen
    doc.close()
    return seiten


def extract_pptx(pptx_path):
    """PPTX: Text pro Folie extrahieren, Foliennummer als Seitennummer."""
    prs = Presentation(pptx_path)
    seiten = []
    for i, folie in enumerate(prs.slides):
        text = ""
        for shape in folie.shapes:
            if shape.has_text_frame:
                for absatz in shape.text_frame.paragraphs:
                    for run in absatz.runs:
                        text += run.text
                    text += "\n"
        seiten.append((i + 1, text))  # Folien ab 1 zählen
    return seiten


def extract_docx(docx_path):
    """DOCX: keine echte Seitenzahl, gesamter Text mit Seitennummer None."""
    doc = Document(docx_path)
    text = ""
    for absatz in doc.paragraphs:
        text += absatz.text + "\n"
    return [(None, text)]  # Seitennummer None -> keine Seitenangabe


def chunks_toembeddings(chunks):
    """Liste von Chunk-Texten in Embeddings umwandeln."""
    embeddings = embedding_model.encode(chunks)
    return embeddings


def save_to_qdrant(chunk_liste, embeddings, source_title, source_type):
    """
    chunk_liste: Liste von (chunk_text, seitennummer)
    embeddings: zugehörige Vektoren in gleicher Reihenfolge
    """
    client = get_qdrant_client()
    points = []
    for (chunk, seitennummer), embedding in zip(chunk_liste, embeddings):
        payload = {
            "content": chunk,
            "source_title": source_title,
            "source_type": source_type,
        }
        # Seitennummer nur speichern, wenn vorhanden (bei DOCX None)
        if seitennummer is not None:
            payload["page_number"] = seitennummer

        point = models.PointStruct(
            id=str(uuid.uuid4()),
            vector=embedding.tolist(),
            payload=payload,
        )
        points.append(point)

    client.upsert(collection_name="documents", points=points)


def ingest_document(file_path, source_title):
    """
    Erkennt anhand der Dateiendung das Format, extrahiert seitenweise,
    zerlegt in Chunks, erzeugt Embeddings und speichert alles in Qdrant.
    """
    endung = os.path.splitext(file_path)[1].lower()

    if endung == ".pdf":
        seiten = extract_pdf(file_path)
        source_type = "pdf"
    elif endung == ".pptx":
        seiten = extract_pptx(file_path)
        source_type = "pptx"
    elif endung == ".docx":
        seiten = extract_docx(file_path)
        source_type = "docx"
    else:
        raise ValueError(f"Nicht unterstütztes Format: {endung}")

    chunk_liste = _split_seiten(seiten)
    if not chunk_liste:
        print(f"Keine Textinhalte in '{source_title}' gefunden.")
        return

    nur_texte = [c for c, _ in chunk_liste]
    embeddings = chunks_toembeddings(nur_texte)
    save_to_qdrant(chunk_liste, embeddings, source_title, source_type)
    print(f"{len(chunk_liste)} Chunks gespeichert aus '{source_title}' ({source_type})")


# Rückwärtskompatibilität: alte Funktion ruft die neue auf
def ingest_pdf(pdf_path, source_title):
    ingest_document(pdf_path, source_title)